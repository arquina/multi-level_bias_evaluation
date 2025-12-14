#!/usr/bin/env python3
"""
01_build_predictions_consistency_overlay_stats.py

Responsibilities
1) Build center-specific prototypes (class centroids) from mean_feature.
2) For each WSI (sample), compute patch-to-prototype distances and write CSV:
   - distance columns: {center}_{subtype}
   - derived columns: {center}_prediction, {center}_confidence
   - consistency columns: all_consistent, consistent_num, consistent_subtype
   - optional: x, y columns (level-0 coords) for safer overlay alignment
3) Render per-sample consistency grid overlay (rows=variants, cols=models) and save.
4) Export one image per sample into a PPT (one slide per sample).
5) Compute non-consistent patch ratio per (WSI, model, variant) and draw boxplot.

Notes
- "all_consistent" is defined as: all three center predictions exist AND identical.
- "consistent_num" / "consistent_subtype" are computed by majority vote across centers.
"""

import os
import glob
import argparse
from typing import Dict, List, Optional, Tuple
from collections import defaultdict

import numpy as np
import pandas as pd
import torch
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.gridspec import GridSpec
from tqdm import tqdm
import openslide
from PIL import Image, ImageDraw, ImageFont

from pptx import Presentation
from pptx.util import Inches


# ---------------------------
# Configuration defaults
# ---------------------------
DEFAULT_VARIANTS = ["original", "stainnorm", "canceronly", "canceronly_stainnorm"]
DEFAULT_MODELS   = ["virchow", "virchow2", "uni_v1", "uni_v2", "gigapath", "conch_v1"]
DEFAULT_PATCH_SIZE = {"virchow":"224", "uni_v1":"256", "gigapath":"256", "virchow2":"224", "uni_v2":"256", "conch_v1":"512"}

# Consistency is computed across these centers (must match column name prefixes)
DEFAULT_CENTERS = ["MSKCC", "NCI Urologic Oncology Branch", "MD Anderson Cancer Center"]
DEFAULT_RACES = ['asian', 'black or african american', 'white']


# ---------------------------
# Utilities: paths & IO
# ---------------------------
def ensure_dir(p: str):
    os.makedirs(p, exist_ok=True)

def clean_patient_name(s: str) -> str:
    # Remove all whitespace for filename matching consistency
    return "".join(str(s).split())

def find_svs_path(svs_root_dir: str, sample: str, subtype: str) -> Optional[str]:
    cand_dir = os.path.join(svs_root_dir, "TCGA", subtype)
    if not os.path.isdir(cand_dir):
        return None
    hits = glob.glob(os.path.join(cand_dir, f"*{sample}*.svs"))
    return hits[0] if hits else None

def get_thumbnail_and_scale(svs_path: str, max_width: int = 2400) -> Tuple[Image.Image, float, float]:
    slide = openslide.OpenSlide(svs_path)
    w0, h0 = slide.dimensions
    target_w = min(max_width, w0)
    scale = target_w / w0
    target_h = int(h0 * scale)
    thumb = slide.get_thumbnail((target_w, target_h)).convert("RGB")
    slide.close()
    return thumb, scale, scale

def load_coords_for_sample(coord_root_dir: str, model: str, patch_size: str, subtype: str, sample: str) -> Optional[np.ndarray]:
    cdir = os.path.join(coord_root_dir, model, patch_size, subtype, "coords")
    if not os.path.isdir(cdir):
        return None
    hits = [p for p in os.listdir(cdir) if (sample in p and p.endswith(".npy"))]
    if not hits:
        return None
    return np.load(os.path.join(cdir, hits[0]))

def find_feature_file(dataset_root: str, feature_type: str, patch_size: str, subtype: str, sample: str) -> Optional[str]:
    d = os.path.join(dataset_root, "whole_cancer", feature_type, patch_size, subtype, "total_feature")
    if not os.path.isdir(d):
        return None
    hits = [p for p in os.listdir(d) if (sample in p and p.endswith((".pt", ".pth")))]
    return os.path.join(d, hits[0]) if hits else None

def find_mean_feature_file(dataset_root: str, feature_type: str, patch_size: str, subtype: str, sample: str, variant: str) -> Optional[str]:
    if 'canceronly' in variant:
        d = os.path.join(dataset_root, "cancer_only", feature_type, patch_size, subtype, "mean_feature")
    else:
        d = os.path.join(dataset_root, "whole_cancer", feature_type, patch_size, subtype, "mean_feature")
    if not os.path.isdir(d):
        return None
    hits = [p for p in os.listdir(d) if (sample in p and p.endswith((".pt", ".pth")))]
    return os.path.join(d, hits[0]) if hits else None


# ---------------------------
# Core: center-wise prediction from distances
# ---------------------------
def save_result_in_dataframe(df: pd.DataFrame) -> pd.DataFrame:
    """
    Input df: distance columns named like "{center}_{subtype}".
    Adds:
      - "{center}_prediction": subtype with minimum distance
      - "{center}_confidence": (d2 - d1) / max(d1, eps) using 1st and 2nd smallest distances
    """
    allowed_subtypes = {"KIRC", "KIRP", "KICH", "normal"}

    center_to_cols = defaultdict(list)
    for col in df.columns:
        if "_" not in col:
            continue
        center, subtype = col.rsplit("_", 1)
        if subtype in allowed_subtypes:
            center_to_cols[center].append(col)

    for center, cols in center_to_cols.items():
        if not cols:
            continue
        df[cols] = df[cols].apply(pd.to_numeric, errors="coerce")

        min_col = df[cols].idxmin(axis=1)
        df[f"{center}_prediction"] = min_col.str.rsplit("_").str[1]

        vals = df[cols].to_numpy(dtype=float)
        d1 = np.nanmin(vals, axis=1)

        col_to_idx = {col: i for i, col in enumerate(cols)}
        argmin_idx = min_col.map(col_to_idx).to_numpy()

        v2 = vals.copy()
        valid_argmin = ~pd.isna(argmin_idx)
        if valid_argmin.any():
            v2[valid_argmin, argmin_idx[valid_argmin].astype(int)] = np.inf

        d2 = np.nanmin(v2, axis=1)
        finite_counts = np.isfinite(vals).sum(axis=1)
        valid_two = finite_counts >= 2

        denom = np.maximum(d1, 1e-6)
        conf = np.full(len(df), np.nan, dtype=float)
        mask = valid_two & np.isfinite(d1) & np.isfinite(d2)
        conf[mask] = (d2[mask] - d1[mask]) / denom[mask]

        df[f"{center}_confidence"] = conf

    return df


# ---------------------------
# Core: patch-level consistency
# ---------------------------
def compute_consistency(df: pd.DataFrame, centers: List[str]) -> pd.DataFrame:
    """
    Adds:
      - all_consistent: all center predictions exist AND identical
      - consistent_num: maximum vote count across centers
      - consistent_subtype: subtype with maximum vote count
    """
    pred_cols = [f"{c}_prediction" for c in centers if f"{c}_prediction" in df.columns]
    if len(pred_cols) == 0:
        df["all_consistent"] = False
        df["consistent_num"] = 0
        df["consistent_subtype"] = np.nan
        return df

    sub = df[pred_cols].copy()

    df["all_consistent"] = (
        sub.notna().all(axis=1) &
        sub.nunique(axis=1, dropna=True).eq(1)
    )

    stacked = sub.stack().dropna()
    if len(stacked) == 0:
        df["consistent_num"] = 0
        df["consistent_subtype"] = np.nan
        return df

    vc = stacked.groupby(level=0).value_counts()
    df["consistent_num"] = vc.groupby(level=0).max().reindex(df.index, fill_value=0).astype(int)

    top_idx = vc.groupby(level=0).idxmax()
    df["consistent_subtype"] = top_idx.map(lambda t: t[1]).reindex(df.index)

    return df


# ---------------------------
# Prototypes
# ---------------------------
def build_center_prototypes(
    meta_df: pd.DataFrame,
    centers: List[str],
    dataset_root: str,
    feature_type: str,
    patch_size: str,
    variant: str,
    subtype_list: List[str] = ["KIRC", "KIRP", "KICH"],
    include_normal: bool = True,
) -> Dict[str, Dict[str, torch.Tensor]]:
    """
    Build prototypes: prototype[center][subtype] = mean over patient mean_feature tensors.
    For 'normal', uses mean_feature under non-cancer path if include_normal=True.
    """
    proto = {c: {s: [] for s in subtype_list} for c in centers}
    if include_normal:
        for c in centers:
            proto[c]["normal"] = []

    # Patient list per center, across the whole metadata subset
    for _, row in meta_df.iterrows():
        patient = clean_patient_name(row["Patient"])
        center  = str(row["center"])
        subtype = str(row["subtype"])
        if center not in centers:
            continue
        if subtype not in subtype_list:
            continue

        mf = find_mean_feature_file(dataset_root, feature_type, patch_size, subtype, patient, variant)
        if mf is not None:
            proto[center][subtype].append(torch.load(mf, weights_only=True))

        if include_normal:
            # normal mean_feature is stored under non-cancer/{target_subtype}/mean_feature in your earlier code
            nd = os.path.join(dataset_root, "non_cancer", feature_type, patch_size, subtype, "mean_feature")
            if os.path.isdir(nd):
                hits = [p for p in os.listdir(nd) if (patient in p and p.endswith((".pt", ".pth")))]
                if hits:
                    proto[center]["normal"].append(torch.load(os.path.join(nd, hits[0]), weights_only=True))

    # Reduce lists -> centroids
    out = {}
    for c in centers:
        out[c] = {}
        for s in subtype_list:
            if len(proto[c][s]) == 0:
                raise RuntimeError(f"No mean_feature found for center={c}, subtype={s}, feature_type={feature_type}, patch_size={patch_size}")
            out[c][s] = torch.mean(torch.stack(proto[c][s], dim=0), dim=0)
        if include_normal:
            if len(proto[c]["normal"]) == 0:
                raise RuntimeError(f"No normal mean_feature found for center={c}, feature_type={feature_type}, patch_size={patch_size}")
            out[c]["normal"] = torch.mean(torch.stack(proto[c]["normal"], dim=0), dim=0)
    return out


def build_prototype_tensor_for_sample(
    meta_df: pd.DataFrame,
    sample_row: pd.Series,
    centers: List[str],
    global_proto: Dict[str, Dict[str, torch.Tensor]],
    dataset_root: str,
    feature_type: str,
    patch_size: str,
    variant: str,
    subtype_list: List[str] = ["KIRC", "KIRP", "KICH"],
    include_normal: bool = True,
) -> Tuple[torch.Tensor, List[str]]:
    """
    Build prototype tensor stacked in deterministic column order:
      center1_KIRC, center1_KIRP, center1_KICH, [center1_normal], center2_..., ...

    For the sample's own center, uses leave-one-out prototype (excluding that sample patient),
    matching your original logic.
    """
    patient = clean_patient_name(sample_row["Patient"])
    sample_center = str(sample_row["center"])

    # Leave-one-out for own center
    loo_proto = None
    if sample_center in centers:
        df_center = meta_df[(meta_df["center"] == sample_center) & (meta_df["Patient"] != patient)]
        # build center-only prototypes from df_center
        tmp_meta = df_center.reset_index(drop=True)
        loo_proto = build_center_prototypes(
            meta_df=tmp_meta,
            centers=[sample_center],
            dataset_root=dataset_root,
            feature_type=feature_type,
            patch_size=patch_size,
            variant = variant,
            subtype_list=subtype_list,
            include_normal=include_normal
        )[sample_center]

    blocks = []
    col_order = []
    for c in centers:
        if c == sample_center and loo_proto is not None:
            center_block = [loo_proto[s] for s in subtype_list]
            if include_normal:
                center_block.append(loo_proto["normal"])
        else:
            center_block = [global_proto[c][s] for s in subtype_list]
            if include_normal:
                center_block.append(global_proto[c]["normal"])

        blocks.append(torch.stack(center_block, dim=0))
        for s in subtype_list:
            col_order.append(f"{c}_{s}")
        if include_normal:
            col_order.append(f"{c}_normal")

    proto_tensor = torch.cat(blocks, dim=0)  # [num_centers * (3 or 4), D]
    return proto_tensor, col_order


# ---------------------------
# Overlay rendering
# ---------------------------
def df_has_xy(df: pd.DataFrame) -> bool:
    return ("x" in df.columns) and ("y" in df.columns)

def draw_consistency_overlay_binary(
    thumb: Image.Image,
    coords: Optional[np.ndarray],
    df: pd.DataFrame,
    patch_size_level0: int,
    scale_x: float,
    scale_y: float,
    color_consistent: Tuple[int,int,int,int] = (0, 102, 204, 100),   # blue
    color_inconsistent: Tuple[int,int,int,int] = (255, 0, 0, 100) # red
) -> Tuple[Image.Image, float]:
    """
    Visualize patch-level consistency as binary colors.
    - all_consistent == True  -> blue
    - all_consistent == False -> red

    Returns
    -------
    overlay_img : PIL.Image (RGB)
    consistent_ratio : float
    """

    overlay = Image.new("RGBA", thumb.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay, "RGBA")

    # --- coordinates ---
    if df_has_xy(df):
        xy = df[["x", "y"]].to_numpy(dtype=float)
    else:
        if coords is None:
            return thumb, float("nan")
        n = min(len(coords), len(df))
        xy = coords[:n]
        df = df.iloc[:n].copy()

    if "all_consistent" not in df.columns:
        return thumb, float("nan")

    ac = df["all_consistent"].fillna(False).to_numpy(dtype=bool)

    ps_w = max(1, int(patch_size_level0 * scale_x))
    ps_h = max(1, int(patch_size_level0 * scale_y))

    for i, (x0, y0) in enumerate(xy):
        X = int(x0 * scale_x)
        Y = int(y0 * scale_y)
        col = color_consistent if ac[i] else color_inconsistent
        draw.rectangle(
            [X, Y, X + ps_w, Y + ps_h],
            fill=col,
            outline=None
        )

    out = thumb.convert("RGBA")
    out.alpha_composite(overlay)

    consistent_ratio = float(ac.mean()) if len(ac) else float("nan")
    return out.convert("RGB"), consistent_ratio


def build_consistency_grid_figure(
    sample: str,
    meta_row: pd.Series,
    variants: List[str],
    models: List[str],
    patch_size_map: Dict[str,str],
    coord_root_dir: str,
    svs_root_dir: str,
    interpret_root: str,
    out_png: str,
    patch_size_multiplier: int = 2,
    max_thumb_width: int = 2400,
) -> str:
    """
    Build a grid:
      rows = variants
      cols = models
    Each cell overlays consistency on a shared WSI thumbnail for the sample.
    """
    subtype = str(meta_row["subtype"])
    svs_path = find_svs_path(svs_root_dir, sample, subtype)
    if svs_path is None:
        return ""

    thumb, sx, sy = get_thumbnail_and_scale(svs_path, max_width=max_thumb_width)

    nrows, ncols = len(variants), len(models)
    fig = plt.figure(figsize=(max(10, 3.2*ncols), max(8, 3.2*nrows)), dpi=150)
    gs = GridSpec(nrows=nrows, ncols=ncols, figure=fig)

    for r, v in enumerate(variants):
        for c, m in enumerate(models):
            ax = fig.add_subplot(gs[r, c])
            csv_path = os.path.join(interpret_root, v, "sample_result", m, f"{sample}.csv")
            if not os.path.exists(csv_path):
                ax.axis("off")
                ax.set_title(f"{v} / {m}\n(no csv)", fontsize=9)
                continue

            df = pd.read_csv(csv_path)
            coords = None
            if not df_has_xy(df):
                coords = load_coords_for_sample(coord_root_dir, m, patch_size_map[m], subtype, sample)

            ps0 = int(patch_size_map[m]) * patch_size_multiplier

            overlay_img, ratio = draw_consistency_overlay_binary(
                thumb=thumb,
                coords=coords,
                df=df,
                patch_size_level0=ps0,
                scale_x=sx,
                scale_y=sy,
                color_inconsistent=(255,0,0,180),
            )

            ax.imshow(overlay_img)
            ax.axis("off")
            ratio_str = "nan" if (ratio != ratio) else f"{ratio:.1%}"
            if r == 0:
                ax.set_title(f"{m}\nconsistent={ratio_str}", fontsize=10)
            else:
                ax.set_title(f"{v}\nconsistent={ratio_str}", fontsize=10)

    fig.suptitle(f"{sample}  |  subtype={subtype}  |  real_center={str(meta_row['center'])}", fontsize=12)
    fig.tight_layout(rect=[0.01, 0.01, 1, 0.95])
    fig.savefig(out_png, bbox_inches="tight")
    plt.close(fig)
    return out_png


def add_slide_one_image(prs: Presentation, image_path: str, margin_in: float = 0.35, top_in: float = 0.35):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # typically blank
    margin = Inches(margin_in)
    top = Inches(top_in)
    width = prs.slide_width - margin*2
    slide.shapes.add_picture(image_path, margin, top, width=width)


# ---------------------------
# Stats: non-consistent ratio
# ---------------------------
def non_consistent_ratio(df: pd.DataFrame) -> float:
    if "all_consistent" not in df.columns:
        return float("nan")
    return float((~df["all_consistent"].fillna(False)).mean())


def collect_non_consistent_stats(
    interpret_root: str,
    variants: List[str],
    models: List[str]
) -> pd.DataFrame:
    rows = []
    for v in variants:
        for m in models:
            d = os.path.join(interpret_root, v, "sample_result", m)
            if not os.path.isdir(d):
                continue
            for f in os.listdir(d):
                if not f.endswith(".csv"):
                    continue
                sample = os.path.splitext(f)[0]
                df = pd.read_csv(os.path.join(d, f))
                rows.append({
                    "sample": sample,
                    "variant": v,
                    "model": m,
                    "non_consistent_ratio": non_consistent_ratio(df),
                })
    return pd.DataFrame(rows)


def plot_non_consistent_boxplot(df: pd.DataFrame, variants: List[str], models: List[str], out_png: str):
    """
    Matplotlib boxplot grouped by model, colored by variant.
    """
    ensure_dir(os.path.dirname(out_png))

    # Prepare data: for each model, variant -> list of ratios
    data = {(m, v): [] for m in models for v in variants}
    for _, r in df.iterrows():
        if r["model"] in models and r["variant"] in variants:
            data[(r["model"], r["variant"])].append(r["non_consistent_ratio"])

    fig = plt.figure(figsize=(max(10, 1.6*len(models)), 6), dpi=160)
    ax = plt.gca()

    group_w = 0.8
    nV = len(variants)
    offsets = np.linspace(-group_w/2, group_w/2, nV, endpoint=False) + (group_w/nV)/2

    positions = []
    box_data = []
    colors = ["#4c78a8", "#f58518", "#54a24b", "#e45756"]  # stable palette for variants
    color_map = {v: colors[i % len(colors)] for i, v in enumerate(variants)}

    for i, m in enumerate(models):
        base = i + 1
        for j, v in enumerate(variants):
            positions.append(base + offsets[j])
            box_data.append(data[(m, v)])

    bp = ax.boxplot(
        box_data,
        positions=positions,
        widths=(group_w/nV)*0.9,
        patch_artist=True,
        showfliers=False,
    )

    for k, patch in enumerate(bp["boxes"]):
        v = variants[k % nV]
        patch.set_facecolor(color_map[v])
        patch.set_alpha(0.7)

    ax.set_xticks(range(1, len(models)+1))
    ax.set_xticklabels(models, rotation=0)
    ax.set_ylabel("Non-consistent patch ratio")
    ax.set_ylim(0, 1)
    ax.grid(axis="y", alpha=0.3, linestyle="--")

    handles = [plt.Line2D([0],[0], color=color_map[v], lw=10, alpha=0.7) for v in variants]
    ax.legend(handles, variants, title="variant", frameon=False, loc="upper right")

    plt.tight_layout()
    plt.savefig(out_png, bbox_inches="tight")
    plt.close(fig)


# ---------------------------
# Main pipeline
# ---------------------------
def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--metadata_csv", help="Path to metadata.csv", default = "/mnt/disk1/Kidney/Submission_dir/metadata.csv")
    parser.add_argument("--svs_root_dir", help="Root for svs data (contains TCGA/{subtype})", default = '/mnt/disk3/svs_data/')
    parser.add_argument("--dataset_root", help="Root for Prototype dataset (contains whole_cancer/...)", default = '/mnt/disk1/Kidney/Submission_dir/dataset/20x/Prototype')
    parser.add_argument("--coord_root_dir", help="Root for coords (model/patch_size/subtype/coords)", default = "/mnt/disk1/Kidney/Submission_dir/dataset/20x/Prototype/whole_cancer/")
    parser.add_argument("--interpret_root", help="Output root (interpretation)", default = "/mnt/disk1/Kidney/final_subission_dir/figure5/")
    parser.add_argument("--variants", nargs="+", default=DEFAULT_VARIANTS)
    parser.add_argument("--models", nargs="+", default=DEFAULT_MODELS)
    parser.add_argument("--centers", nargs="+", default=DEFAULT_CENTERS)
    parser.add_argument("--include_normal", default = True, action="store_true", help="Include normal prototype and normal distance columns")
    parser.add_argument("--write_xy", default = True, action="store_true", help="Write x,y columns into CSV for robust overlay alignment")
    parser.add_argument("--patch_size_multiplier", type=int, default=2)
    parser.add_argument("--max_thumb_width", type=int, default=2400)
    parser.add_argument("--make_ppt", action="store_true", default = True)
    parser.add_argument("--skip_csv", action="store_true", help="Skip CSV generation (only visualize/stats)", default = True)
    parser.add_argument("--skip_overlay", action="store_true", help="Skip overlay generation/PPT", default = False)
    parser.add_argument("--skip_stats", action="store_true", help="Skip stats/boxplot", default = False)
    args = parser.parse_args()
    
    # patch size map
    patch_size_map = dict(DEFAULT_PATCH_SIZE)
    ensure_dir(args.interpret_root)
    meta = pd.read_csv(args.metadata_csv)
    meta = meta.drop_duplicates(subset=["Patient"]).copy()
    # normalize patient names for matching
    meta["Patient"] = meta["Patient"].astype(str).map(clean_patient_name)

    # filter to kidney subtypes if present
    meta = meta[meta["subtype"].isin(["KIRC", "KIRP", "KICH"])].reset_index(drop=True)
    meta = meta[meta['center'].isin(DEFAULT_CENTERS)]
    meta = meta[meta['race'].isin(DEFAULT_RACES)]

    # ---- 1) CSV generation ----
    if not args.skip_csv:
        print("[STEP] CSV generation: distances, predictions, confidence, consistency")

        for variant in args.variants:
            is_stainnorm = ("stainnorm" in variant)
            feature_suffix = "_stainnorm" if is_stainnorm else ""

            for model in args.models:
                patch_size = str(patch_size_map[model])
                feature_type = f"{model}{feature_suffix}"

                # Build global prototypes once per (variant, model)
                print(f"  - Build prototypes: variant={variant}, model={model}, feature_type={feature_type}, patch_size={patch_size}")
                proto_global = build_center_prototypes(
                    meta_df=meta,
                    centers=args.centers,
                    dataset_root=args.dataset_root,
                    feature_type=feature_type,
                    patch_size=patch_size,
                    variant = variant,
                    include_normal=args.include_normal,
                )

                out_dir = os.path.join(args.interpret_root, variant, "sample_result", model)
                ensure_dir(out_dir)

                # Iterate samples
                for idx in tqdm(range(len(meta)), desc=f"CSV {variant}/{model}", leave=False):
                    row = meta.iloc[idx]
                    sample = row["Patient"]
                    subtype = str(row["subtype"])

                    # Load features
                    feat_path = find_feature_file(args.dataset_root, feature_type, patch_size, subtype, sample)
                    if feat_path is None:
                        continue
                    feats = torch.load(feat_path, weights_only=True)
                    if isinstance(feats, np.ndarray):
                        feats = torch.from_numpy(feats)
                    feats = feats.float()

                    # Build prototype tensor (with leave-one-out for own center)
                    proto_tensor, col_order = build_prototype_tensor_for_sample(
                        meta_df=meta,
                        sample_row=row,
                        centers=args.centers,
                        global_proto=proto_global,
                        dataset_root=args.dataset_root,
                        feature_type=feature_type,
                        patch_size=patch_size,
                        variant = variant,
                        include_normal=args.include_normal,
                    )

                    # Distances
                    dist = torch.cdist(feats, proto_tensor).cpu().numpy()
                    df = pd.DataFrame(dist, columns=col_order)

                    # predictions + confidence
                    df = save_result_in_dataframe(df)

                    # consistency
                    df = compute_consistency(df, centers=args.centers)

                    # optionally attach x,y
                    if args.write_xy:
                        coords = load_coords_for_sample(args.coord_root_dir, model, patch_size, subtype, sample)
                        if coords is not None and len(coords) >= len(df):
                            df["x"] = coords[:len(df), 0].astype(int)
                            df["y"] = coords[:len(df), 1].astype(int)

                    df.to_csv(os.path.join(out_dir, f"{sample}.csv"), index=False)

    # ---- 2) Overlay grid figures + PPT ----
    grid_dir = os.path.join(args.interpret_root, "visualization", "consistency_grid")
    ensure_dir(grid_dir)

    ppt_path = os.path.join(args.interpret_root, "figure", "consistency_grids__one_per_sample.pptx")
    ensure_dir(os.path.dirname(ppt_path))

    image_paths = {}

    if not args.skip_overlay:
        print("[STEP] Overlay grid figures")
        for idx in tqdm(range(len(meta)), desc="Overlay grids"):
            row = meta.iloc[idx]
            sample = row["Patient"]

            out_png = os.path.join(grid_dir, f"{sample}__CONSISTENCY_GRID.png")
            img = build_consistency_grid_figure(
                sample=sample,
                meta_row=row,
                variants=args.variants,
                models=args.models,
                patch_size_map=patch_size_map,
                coord_root_dir=args.coord_root_dir,
                svs_root_dir=args.svs_root_dir,
                interpret_root=args.interpret_root,
                out_png=out_png,
                patch_size_multiplier=args.patch_size_multiplier,
                max_thumb_width=args.max_thumb_width,
            )
            if img:
                image_paths[sample] = img

        if args.make_ppt:
            print("[STEP] PPT export")
            prs = Presentation()
            for sample in tqdm(list(image_paths.keys()), desc="PPT slides"):
                add_slide_one_image(prs, image_paths[sample], margin_in=0.35, top_in=0.35)
            prs.save(ppt_path)
            print(f"[DONE] PPT saved -> {ppt_path}")

    # ---- 3) Stats / boxplot ----
    if not args.skip_stats:
        print("[STEP] Non-consistent ratio stats + boxplot")
        stats_df = collect_non_consistent_stats(args.interpret_root, args.variants, args.models)

        stats_dir = os.path.join(args.interpret_root, "stats")
        ensure_dir(stats_dir)

        stats_csv = os.path.join(stats_dir, "non_consistent_ratios.csv")
        stats_df.to_csv(stats_csv, index=False)

        out_png = os.path.join(stats_dir, "boxplot_non_consistent_ratio.png")
        plot_non_consistent_boxplot(stats_df, args.variants, args.models, out_png)

        print(f"[DONE] stats csv -> {stats_csv}")
        print(f"[DONE] boxplot   -> {out_png}")


if __name__ == "__main__":
    main()
